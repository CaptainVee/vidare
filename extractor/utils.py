import PyPDF2
from pptx import Presentation
from pdfminer.high_level import extract_text
import pdfminer
import pdfminer.settings

pdfminer.settings.STRICT = False  # Disable strict mode for more lenient parsing

from pdfminer.pdfparser import PDFParser
from pdfminer.pdfdocument import PDFDocument


def parse_pdf(pdf_path):
    document_info = {}

    with open(pdf_path, "rb") as file:
        parser = PDFParser(file)
        document = PDFDocument(parser)
        pdf_reader = PyPDF2.PdfReader(file)
        num_pages = len(pdf_reader.pages)

        document_info["Content"] = extract_text(file)
        document_info["no_of_pages"] = num_pages

        for info_dict in document.info:
            for key, value in info_dict.items():
                document_info[key] = value.decode()

    return document_info


def parse_pptx(file_path):
    extracted_text = ""
    no_of_pages = 0
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        no_of_pages += 1
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    extracted_text += paragraph.text + "\n"

    return {"Content": extracted_text, "no_of_pages": no_of_pages}
